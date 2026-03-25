#!/usr/bin/env python3
"""Generate PPTX from leadership presentation data."""

import sys
sys.path.insert(0, '/tmp/pdfenv/lib/python3.14/site-packages')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(0x0f, 0x11, 0x17)
WHITE = RGBColor(0xf1, 0xf5, 0xf9)
MUTED = RGBColor(0x94, 0xa3, 0xb8)
BLUE = RGBColor(0x4f, 0x8c, 0xff)
GREEN = RGBColor(0x34, 0xd3, 0x99)
AMBER = RGBColor(0xfb, 0xbf, 0x24)
RED = RGBColor(0xf8, 0x71, 0x71)
PURPLE = RGBColor(0xa7, 0x8b, 0xfa)
CYAN = RGBColor(0x22, 0xd3, 0xee)
DARK_CARD = RGBColor(0x1a, 0x1d, 0x27)

def add_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG

def add_title_slide(title, subtitle_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_bg(slide)
    t = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(20)
    p2.font.color.rgb = MUTED
    p2.alignment = PP_ALIGN.CENTER
    return slide

def add_content_slide(title, bullets, color=BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    t = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
    tf = t.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color

    body = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.5), Inches(5.5))
    tf2 = body.text_frame
    tf2.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE if not bullet.startswith("  ") else MUTED
        p.space_after = Pt(6)
    return slide

# SLIDE 1: Title
s = add_title_slide("Sportsbook Programme", "World Cup MVP Delivery Plan | March 2026")
info = s.shapes.add_textbox(Inches(2), Inches(4), Inches(9), Inches(1))
tf = info.text_frame
p = tf.paragraphs[0]
p.text = "16 Work Packages  |  4 Deliverables  |  8 Domain Teams"
p.font.size = Pt(18)
p.font.color.rgb = BLUE
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "Soft Launch: Apr 1  |  Go-Live: May 1  |  World Cup: Mid-June 2026"
p2.font.size = Pt(14)
p2.font.color.rgb = MUTED
p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph()
p3.text = "Scope: Soccer only | Pre-match | USD | en-US | Web (desktop + mobile) | 4 brands"
p3.font.size = Pt(12)
p3.font.color.rgb = MUTED
p3.alignment = PP_ALIGN.CENTER

# SLIDE 2: 4 Deliverables
add_content_slide("4 Deliverables — Each a Demoable Increment", [
    "1. Central Offering Foundation",
    "  Punter opens the sportsbook and browses World Cup fixtures with live odds",
    "  6 WPs: FEEDS, OFFER, FRONTEND, BETTING, DATA",
    "",
    "2. Betslip Core & Placement Flow",
    "  Punter builds a bet, places it, gets validated. Operator controls limits and kill switch",
    "  11 WPs: FRONTEND, BETTING, OFFER, DATA",
    "",
    "3. My Bets & Real Data Connection",
    "  Bets settle from real results. Real money moves. Punter sees bet history live",
    "  10 WPs: BETTING, DATA, FRONTEND, FEEDS, OFFER",
    "",
    "4. Operation Control & Visibility",
    "  Traders manage the book via Kraken. Splunk dashboards show full operational picture",
    "  14 WPs: BACKOFFICE, BETTING, DATA, SPLUNK",
])

# SLIDE 3: 16 Work Packages
add_content_slide("16 Work Packages", [
    "1  Simple Parlay — Combine bets from different matches",
    "2  Soccer WC Hub — World Cup landing page with fixtures & markets",
    "3  Settlement — Auto-grade and settle bets from real results",
    "4  Book On/Off Switch — Emergency kill switch + punter blocking",
    "5  Bet Placement Controls — Stake limits, payout caps, odds validation",
    "6  Central Offer Foundation — Own our offer data, not the feed provider's",
    "7  Bet Ticker — Live bet feed for traders",
    "8  RT Liability — Real-time book exposure by market",
    "9  Manual Resettlement — Traders correct grading errors and void bets",
    "10 Audit Trail — Who changed what, when, why",
    "11 RT Odds Display — Live odds updates + suspended market indicators",
    "12 My Bets — Punter checks active and settled bets",
    "13 Bet History — Traders investigate punter betting activity",
    "14 SGP — Same-game parlay with correlated pricing",
    "15 Mega Parlay — SGPs + straights across matches",
    "16 Feed Health — Detect unhealthy feeds, protect from stale data",
])

# SLIDE 4: Jira Status
add_content_slide("Work Package Status — Live from Jira (25 Mar 2026)", [
    "1 COMPLETED: Straight Bet (SA-15) — Post-Delivery Assessment",
    "",
    "11 IN DELIVERY:",
    "  SA-66 Central Offer | SA-14 WC Hub | SA-17 Simple Parlay | SA-22 Settlement",
    "  SA-20 Bet Controls | SA-25 Bet Ticker | SA-24 RT Liability | SA-26 Bet History",
    "  SA-27 My Bets | SA-71 Feed Health | SA-28 Audit Trail",
    "",
    "4 READY FOR HANDOVER:",
    "  SA-19 RT Odds Display | SA-21 Book On/Off | SA-23 Manual Resettlement | SA-28 Audit Trail",
    "",
    "2 SCOPE VALIDATION (blocked):",
    "  SA-16 SGP | SA-18 Mega Parlay",
], GREEN)

# SLIDES 5-8: Deliverables
for d_num, d_name, d_color, teams in [
    (1, "Central Offering Foundation", BLUE, [
        "FEEDS — WP 2,6,11,16 — WC data flows from OpticOdds. No feed = no sportsbook",
        "OFFER — WP 2,6,11,16 — We own our offer, not a pass-through of OpticOdds",
        "FRONTEND — WP 2,6,11,16 — WC Hub page, navigation, real-time odds, suspended markets",
        "BETTING — WP 11,16 — SGP prices refresh on odds change. Stop bets when feed is down",
        "DATA — WP 2,6,11,16 — Odds history in Splunk. Diffusion for real-time frontend",
        "SPLUNK RAW — WP 2,6,16 — Offer events, odds history, feed health dashboard",
    ]),
    (2, "Betslip Core & Placement Flow", GREEN, [
        "FRONTEND — WP 1,4,5,14,15 — Combined odds, winnings, errors. Kill switch darkens site. SGP feedback",
        "BETTING — WP 1,4,5,7,8,13,14,15,16 — Core upgrade: parlay+SGP+Mega. Limits. Kill switch. Real wallet",
        "OFFER — WP 4,5,14,15 — Business defines rules: limits, bet types, SGP/Mega config",
        "DATA — WP 1,4,5,7,8,13,14,15,16 — Every bet in reports. Ticker+Liability start. Rejected bets",
        "SPLUNK RAW — WP 1,16 — First multi-leg bet events. Feed health monitoring",
    ]),
    (3, "My Bets & Real Data Connection", PURPLE, [
        "FRONTEND — WP 3,12 — Punter sees Pending>Won/Lost. My Bets = core punter experience",
        "BETTING — WP 3,7,8,9,12,13,14,15 — Where the business makes/loses money. Real grading+wallet",
        "FEEDS — WP 3 — Match results from OpticOdds. Without results, no bets settle",
        "OFFER — WP 3 — Results through our offer resolution. We own the grading",
        "DATA — WP 3,7,8,9,12,13 — Settlement feeds P&L. Liability updates. Manual audit. My Bets push",
        "SPLUNK RAW — WP 3,12 — Every grading decision, wallet transaction, resettlement",
    ]),
    (4, "Operation Control & Visibility", AMBER, [
        "BACKOFFICE — WP 4,5,9,10 — Kraken: kill switch, limits, manual settle, audit trail",
        "BETTING — WP 4,5,7,9,10,13 — Kill switch blocks API. Limits enforced. Manual settle engine. Audit",
        "OFFER — WP 4,5 — Operator Management: master switch and limit config",
        "DATA — WP 4,5,7,8,9,10,11,13,16 — All dashboards: Ticker, Liability, Audit, Odds, Feed Health",
        "SPLUNK BUSINESS — WP 4-11,13-16 — Trader's operational cockpit for game day",
    ]),
]:
    add_content_slide(f"Deliverable {d_num}: {d_name}", teams, d_color)

# SLIDE 9: Team Workload
add_content_slide("Domain Team Workload", [
    "WP assignments per team across all 4 deliverables:",
    "",
    "  Feeds:       D1=4  D2=1  D3=1  D4=1  Total=7",
    "  Offer:       D1=4  D2=4  D3=1  D4=2  Total=11",
    "  Frontend:    D1=4  D2=5  D3=2  D4=—  Total=11",
    "  Betting:     D1=2  D2=9  D3=8  D4=6  Total=25  *** HEAVIEST",
    "  Data:        D1=4  D2=9  D3=6  D4=9  Total=28  *** HEAVIEST",
    "  BackOffice:  D1=—  D2=—  D3=—  D4=4  Total=4",
    "  Splunk:      D1=3  D2=2  D3=2  D4=12 Total=19",
    "",
    "Betting and Data carry the heaviest load across the programme.",
])

# SLIDE 10: Dependencies & Risks
add_content_slide("Dependencies & Risks", [
    "CRITICAL PATH:",
    "  WP 6 Central Offer > WP 2, 11, and all downstream (Foundation)",
    "  WES Wallet > WP 1, 3, 5, 9 (Blocks 4 WPs — real money can't move)",
    "  PAM Auth > WP 1, 5 (Blocks 2 WPs)",
    "  WP 1 Simple Parlay > WP 14 SGP > WP 15 Mega Parlay (Sequential)",
    "",
    "TOP RISKS:",
    "  [HIGH] Mega Parlay & SGP still in Scope Validation — cannot estimate or plan",
    "  [HIGH] Shared DB schema refactor (SBK-2008) — blocks Simple Parlay + Settlement",
    "  [HIGH] WES Wallet integration timeline — cascades across 4 WPs",
    "  [MED]  Splunk vs New Relic strategy undefined — affects all reporting WPs",
    "  [MED]  Feed Health > Betting suspension unscoped — no WP owns this end-to-end",
], RED)

# SLIDE 11: Scope
add_content_slide("Scope: In vs Out (v1.55)", [
    "IN SCOPE (MUST HAVE):",
    "  Pre-match betting on World Cup fixtures",
    "  Straight, Simple Parlay, SGP, Mega Parlay",
    "  Real settlement from OpticOdds results",
    "  Operator controls (kill switch, limits, manual settle)",
    "  Bet Ticker, Liability, Audit Trail in Splunk",
    "  Web desktop + mobile responsive | 4 brands | USD | en-US",
    "  5,000 bets/min | 500ms odds latency | 99.9% uptime",
    "",
    "SHOULD HAVE: In-Play betting | Futures (straights only) | Fractional odds",
    "",
    "OUT OF SCOPE (v1.7+): Multi-currency | Multi-language | Cash Out | Native apps",
    "  Custom markets | Dynamic margins | Per-punter value factors | Other sports",
    "",
    "REDUCED FOR v1.55: Traders Hub UI = API only | BackOffice limits absorbed into Betting",
])

# SLIDE 12: Zoom Levels
add_content_slide("Planning Framework — 4 Zoom Levels", [
    "ZOOM 01 — Business Value Mapping [DONE]",
    "  Gate: Product + WPs defined",
    "  What to demo, what work per domain. Output: 4 Deliverables, 16 WPs, business WHY",
    "",
    "ZOOM 02 — Order & Estimates [NEXT]",
    "  Gate: After Architecture Handover",
    "  Sequence per domain, T-shirt estimates (weeks), cross-domain dependencies",
    "",
    "ZOOM 03 — Milestones & Demo Dates [PENDING]",
    "  Gate: After Software Design",
    "  Concrete dates per deliverable demo, man-day estimates per feature",
    "",
    "ZOOM 04 — Sprint-Level JIRA Plan [PENDING]",
    "  Gate: After Team Refinement",
    "  Story points per sprint, detailed assignments, velocity-based forecasting",
], BLUE)

# SLIDE 13: SDLC Funnel
add_content_slide("SDLC Funnel — How Work Packages Progress", [
    "Product > Architecture > SOFTWARE DESIGN (UNLOCK) > Team Refine > Execution",
    "",
    "Owners per stage:",
    "  Product: Kurt / Stakeholders",
    "  Architecture: Matus / Solutions",
    "  Software Design: Team Principal  <<<< UNLOCK POINT",
    "  Team Refine: Team Refine (estimates)",
    "  Execution: AWS Stage > PPD > Production (full delivery pipeline)",
    "",
    "After Software Design, work fans out into parallel streams:",
    "  Dev Estimates > Sprint > Demo",
    "  QA Estimates > Sprint > Demo",
    "  DevOps Plan > Sprint > Demo",
    "  Data Plan > Sprint > Demo",
    "  FE Plan (can start before SD with mocks)",
    "",
    "Sprint Review 1 > Sprint Review 2 > Deliverable Demo",
], AMBER)

# SLIDE 14: WP Pipeline
add_content_slide("WP Progress Through the Funnel", [
    "Product+Arch done, now in Software Design:",
    "  SA-66 Central Offer | SA-14 WC Hub | SA-17 Simple Parlay",
    "  SA-22 Settlement | SA-20 Bet Controls | SA-25 Bet Ticker",
    "  SA-24 RT Liability | SA-27 My Bets | SA-71 Feed Health",
    "",
    "Architecture complete, waiting for delivery pickup:",
    "  SA-19 RT Odds Display | SA-21 Book On/Off",
    "  SA-23 Manual Resettlement | SA-28 Audit Trail",
    "",
    "Still in Scope Validation (BLOCKED):",
    "  SA-16 SGP | SA-18 Mega Parlay",
    "",
    "Execution = AWS Stage > PPD Done > Production",
    "  Full delivery tracked across 3 environments",
    "",
    "Key: 9 WPs in Software Design. SD completion unlocks parallel estimation.",
    "The SD column is the bottleneck we track weekly.",
])

# SLIDE 15: Epic Progress
add_content_slide("Epic Progress — Child Tickets (SBK-669)", [
    "Epic                              Team            Tickets  Done%",
    "─────────────────────────────────────────────────────────────────",
    "[FE] Straight Bet (SBK-1951)      Frontend          39     33%  (13 done / 13 prog / 13 todo)",
    "[BE] Simple Parlay (SBK-2199)     Betting Engine     34     32%  (11 done / 15 prog / 8 todo)",
    "[BE] Central Offer (SBK-3006)     Offer Mgmt         17      6%  (1 done / 8 prog / 8 todo)",
    "[WC] Market Parser (SBK-1897)     Feeds Mgmt         11     45%  (5 done / 1 prog / 5 todo)",
    "[BE] Feed Health (SBK-3094)       Feeds Integ        10      0%  (0 done / 4 prog / 6 todo)",
    "[BE] Settlement (SBK-2200)        Delivery            9      0%  (0 done / 3 prog / 6 todo)",
    "Soccer WC Hub (SBK-2189)          Delivery            9      0%  (0 done / 0 prog / 9 todo)",
    "Bet Ticker/History (DATA-22)      Data                1      0%  (0 done / 1 prog / 0 todo)",
    "Resulting PreGame (SBK-3516)      Feeds Mgmt          1      0%  (0 done / 0 prog / 1 todo)",
    "[QA] Straight Bet (SBK-3790)      BE QAs              1      0%  (0 done / 1 prog / 0 todo)",
    "",
    "TOTALS: 132 tickets | 30 Done (23%) | 48 In Progress | 54 To Do",
], GREEN)

# SLIDE 16: Next Steps
add_content_slide("Where We Are & Next Steps", [
    "CURRENT POSITION: Zoom Level 01 DONE | Zoom Level 02 STARTING",
    "WP PIPELINE: 1 Done | 11 In Delivery | 4 Ready for Handover | 2 Scope Validation",
    "",
    "IMMEDIATE ACTIONS:",
    "",
    "  1. Unblock SGP & Mega Parlay",
    "     Architecture to complete scope validation for SA-16 and SA-18",
    "",
    "  2. Start Zoom Level 02 for 15 WPs",
    "     DMs fill in delivery order, T-shirt estimates, dependencies",
    "",
    "  3. Platform Decision: Splunk vs New Relic",
    "     Affects 8+ reporting/audit WPs. Decision needed before Z02 estimates for Data",
    "",
    "  4. Confirm External Timelines",
    "     WES Wallet (blocks 4 WPs) and PAM Auth (blocks 2 WPs)",
    "",
    "  5. Pick Up 4 Handover WPs",
    "     RT Odds, Book On/Off, Manual Resettlement, Audit Trail — waiting for capacity",
])

prs.save('/Users/kurtcarabott/WKS-FT-DELIVERY-MANAGEMENT/reports/leadership-presentation.pptx')
print("PPTX saved!")
